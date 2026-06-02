# -*- coding: utf-8 -*-
"""
Presentación final AI-900T00 — Caso InnovVentas / Chatbot Nova.
Estilo minimalista y elegante, LIGHT MODE (exposición en entorno claro).

Hilo narrativo:  DIAGNÓSTICO → CAUSA → SOLUCIÓN → IMPLEMENTACIÓN
(el "kicker" de cada slide indica la fase). Incluye decisión tecnológica
(por qué no Azure y las alternativas elegidas), tecnologías, COSTOS, y las
5 preguntas guía del PDF oficial.

Re-ejecutable:  python entregable/build_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

OUTPUT = "entregable/PRESENTACION-InnovVentas-Nova.pptx"

INK    = RGBColor(0x16, 0x1D, 0x2F)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
LIGHT  = RGBColor(0xF4, 0xF7, 0xFB)
LINE   = RGBColor(0xE2, 0xE8, 0xF0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x05, 0x96, 0x69)
GREENB = RGBColor(0xEC, 0xFD, 0xF5)
RED    = RGBColor(0xDC, 0x26, 0x2A)
REDB   = RGBColor(0xFE, 0xF2, 0xF2)
AMBER  = RGBColor(0xB4, 0x53, 0x09)
AMBERB = RGBColor(0xFF, 0xFB, 0xEB)

F, FL, FSB = "Segoe UI", "Segoe UI Light", "Segoe UI Semibold"
MX = Inches(0.92)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
_page = [0]


def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    return s


def textbox(slide, l, t, w, h, paras, anchor=None):
    box = slide.shapes.add_textbox(l, t, w, h); tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if anchor: tf.vertical_anchor = anchor
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pa.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(pa.get("sa", 4)); p.space_before = Pt(pa.get("sb", 0))
        if "line" in pa: p.line_spacing = pa["line"]
        for rr in pa.get("runs", [pa]):
            r = p.add_run(); r.text = rr.get("text", "")
            r.font.size = Pt(rr.get("size", 16)); r.font.color.rgb = rr.get("color", INK)
            r.font.name = rr.get("font", F); r.font.bold = rr.get("bold", False)
            r.font.italic = rr.get("italic", False)
    return box


def rect(slide, l, t, w, h, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def card(slide, l, t, w, h, fill=LIGHT, border=None, radius=0.06):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if border: sp.line.color.rgb = border; sp.line.width = Pt(1)
    else: sp.line.fill.background()
    sp.shadow.inherit = False
    try: sp.adjustments[0] = radius
    except Exception: pass
    return sp


def arrow(slide, x1, y1, x2, y2, color=ACCENT, w=1.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = Pt(w)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle'}))
    c.shadow.inherit = False
    return c


def header(slide, kicker, title):
    rect(slide, MX, Inches(0.62), Inches(0.30), Inches(0.30), ACCENT)
    textbox(slide, Inches(1.34), Inches(0.60), Inches(10.5), Inches(0.34),
            [{"text": kicker.upper(), "size": 12.5, "color": ACCENT, "font": FSB, "bold": True}])
    textbox(slide, MX, Inches(1.02), Inches(11.8), Inches(0.8),
            [{"text": title, "size": 30, "color": INK, "font": FL}])
    rect(slide, MX, Inches(1.78), Inches(0.62), Pt(3), ACCENT)


def footer(slide):
    _page[0] += 1
    textbox(slide, MX, Inches(7.02), Inches(10), Inches(0.3),
            [{"text": "InnovVentas · Chatbot Nova · AI-900T00 · SENATI", "size": 9, "color": MUTED}])
    textbox(slide, Inches(12.2), Inches(7.02), Inches(0.9), Inches(0.3),
            [{"text": str(_page[0]), "size": 9, "color": MUTED, "align": PP_ALIGN.RIGHT}])


# ===== 1 · PORTADA =====
s = blank()
rect(s, 0, 0, Inches(0.22), Inches(7.5), ACCENT)
textbox(s, MX, Inches(1.65), Inches(11), Inches(0.4),
        [{"text": "TRABAJO FINAL DEL CURSO · AI-900T00", "size": 14, "color": ACCENT, "font": FSB, "bold": True}])
textbox(s, MX, Inches(2.2), Inches(11.6), Inches(2.0),
        [{"text": "Nova", "size": 66, "color": INK, "font": FL},
         {"text": "Chatbot inteligente para el E-commerce de InnovVentas", "size": 25, "color": INK, "font": FL, "sb": 6}])
rect(s, MX, Inches(4.5), Inches(2.2), Pt(2.5), LINE)
textbox(s, MX, Inches(4.75), Inches(11.5), Inches(1.6),
        [{"text": "Conceptos Básicos de IA en Microsoft Azure · Tecnologías de la Información", "size": 14, "color": MUTED, "sa": 10},
         {"runs": [{"text": "Equipo:  ", "size": 14, "color": INK, "font": FSB, "bold": True},
                   {"text": "Hugo André Cahua Solano · Henry Humberto Cruces Castro · Jean Beckan Olivitos Villanueva", "size": 14, "color": MUTED}]}])

# ===== 2 · AGENDA =====
s = blank(); header(s, "Contenido", "Agenda")
items = ["Diagnóstico: contexto y problemas", "Causa raíz y propuesta de valor",
         "La solución: Nova", "Decisión tecnológica y costos",
         "Diseño: FAQs, flujo y arquitectura", "Plan de implementación",
         "Métricas y preguntas guía", "Demo en vivo y conclusiones"]
for col, group in enumerate([items[:4], items[4:]]):
    x = MX if col == 0 else Inches(6.85)
    for i, it in enumerate(group):
        n = col * 4 + i + 1; yy = 2.2 + i * 0.9
        card(s, x, Inches(yy), Inches(5.6), Inches(0.68))
        textbox(s, x + Inches(0.2), Inches(yy), Inches(0.7), Inches(0.68),
                [{"text": f"0{n}", "size": 17, "color": ACCENT, "font": FSB, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + Inches(0.9), Inches(yy), Inches(4.55), Inches(0.68),
                [{"text": it, "size": 13.5, "color": INK}], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ===== 3 · DIAGNÓSTICO: CONTEXTO =====
s = blank(); header(s, "Diagnóstico", "Contexto: InnovVentas")
textbox(s, MX, Inches(2.05), Inches(6.0), Inches(3.6),
        [{"text": "Empresa peruana especializada en la comercialización de productos tecnológicos. Su principal canal de ventas es el E-commerce.", "size": 16, "color": INK, "line": 1.25, "sa": 12},
         {"text": "Presenta una baja interacción en su sitio web que impacta las ventas y la retención: los usuarios abandonan la compra por falta de respuestas rápidas a sus dudas.", "size": 16, "color": MUTED, "line": 1.25}])
needs = [("Inmediatez", "Respuestas en segundos, 24/7"),
         ("Cobertura", "Antes, durante y después de la compra"),
         ("Integración", "No intrusiva, sobre el sitio actual"),
         ("Métricas", "Medir impacto en ventas y abandono")]
textbox(s, Inches(7.15), Inches(1.98), Inches(5), Inches(0.4),
        [{"text": "NECESIDADES DEL CLIENTE", "size": 12, "color": ACCENT, "font": FSB, "bold": True}])
for i, (t, d) in enumerate(needs):
    yy = 2.45 + i * 0.95
    card(s, Inches(7.15), Inches(yy), Inches(5.25), Inches(0.82))
    textbox(s, Inches(7.4), Inches(yy), Inches(4.8), Inches(0.82),
            [{"runs": [{"text": t + " — ", "size": 14, "color": INK, "font": FSB, "bold": True},
                       {"text": d, "size": 13, "color": MUTED}]}], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ===== 4 · DIAGNÓSTICO: PROBLEMAS =====
s = blank(); header(s, "Diagnóstico", "Problemas en el flujo de E-commerce")
probs = [("Carritos abandonados", "Los usuarios dejan la compra a medias por dudas no resueltas a tiempo."),
         ("Consultas sin respuesta", "Especificaciones, métodos de pago y disponibilidad sin atención inmediata."),
         ("Soporte lento", "Tickets que tardan 24–48 h generan insatisfacción y fuga a la competencia."),
         ("Sin visibilidad", "No se miden las interacciones; no se puede optimizar la conversión.")]
for i, (t, d) in enumerate(probs):
    x = MX + (i % 2) * Inches(5.95); y = 2.25 + (i // 2) * 2.05
    card(s, x, Inches(y), Inches(5.6), Inches(1.75), fill=REDB)
    rect(s, x, Inches(y), Inches(0.10), Inches(1.75), RED)
    textbox(s, x + Inches(0.35), Inches(y + 0.22), Inches(5.0), Inches(1.4),
            [{"text": t, "size": 17, "color": INK, "font": FSB, "bold": True, "sa": 6},
             {"text": d, "size": 13, "color": MUTED, "line": 1.2}])
footer(s)

# ===== 5 · CAUSA → SOLUCIÓN (centro del hilo narrativo) =====
s = blank(); header(s, "Causa → Solución", "De la causa raíz a la solución")
triples = [
    ("Carritos abandonados", "No hay respuestas rápidas durante la compra", "Nova responde FAQs al instante y guía el checkout."),
    ("Fuga a la competencia", "Dudas de specs, pago y stock sin atención", "Atención 24/7 con base de FAQs; deriva a humano si hace falta."),
    ("Sin mejora continua", "No se miden las interacciones", "Registro de todo + dashboard (CSAT, intenciones, resolución)."),
]
# encabezados de columna
heads = [("PROBLEMA", RED, MX, Inches(3.3)), ("CAUSA", AMBER, MX + Inches(3.75), Inches(3.3)),
         ("SOLUCIÓN — NOVA", GREEN, MX + Inches(7.5), Inches(4.0))]
for txt, col, x, w in heads:
    textbox(s, x, Inches(2.0), w, Inches(0.3), [{"text": txt, "size": 11, "color": col, "font": FSB, "bold": True}])
y = 2.4
for p, c, sol in triples:
    card(s, MX, Inches(y), Inches(3.3), Inches(1.25), fill=REDB)
    textbox(s, MX + Inches(0.2), Inches(y), Inches(2.95), Inches(1.25),
            [{"text": p, "size": 13, "color": INK, "font": FSB, "bold": True, "line": 1.1}], anchor=MSO_ANCHOR.MIDDLE)
    arrow(s, MX + Inches(3.32), Inches(y + 0.625), MX + Inches(3.73), Inches(y + 0.625), color=MUTED)
    card(s, MX + Inches(3.75), Inches(y), Inches(3.3), Inches(1.25), fill=AMBERB)
    textbox(s, MX + Inches(3.95), Inches(y), Inches(2.95), Inches(1.25),
            [{"text": c, "size": 12.5, "color": INK, "line": 1.12}], anchor=MSO_ANCHOR.MIDDLE)
    arrow(s, MX + Inches(7.07), Inches(y + 0.625), MX + Inches(7.48), Inches(y + 0.625))
    card(s, MX + Inches(7.5), Inches(y), Inches(4.0), Inches(1.25), fill=GREENB)
    textbox(s, MX + Inches(7.7), Inches(y), Inches(3.65), Inches(1.25),
            [{"text": sol, "size": 12.5, "color": INK, "line": 1.12}], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.42
footer(s)

# ===== 6 · SOLUCIÓN: NOVA =====
s = blank(); header(s, "Solución", "Nova, el asistente virtual")
textbox(s, MX, Inches(2.05), Inches(11.4), Inches(1.0),
        [{"text": "Chatbot conversacional embebible en el sitio de InnovVentas. Entiende lenguaje natural con un modelo de lenguaje (LLM) y responde ciñéndose a las FAQs de la empresa; registra cada conversación y mide la satisfacción del cliente.", "size": 16, "color": MUTED, "line": 1.25}])
feats = [("Responde FAQs", "Productos, pagos, envíos, garantías y devoluciones."),
         ("Asiste la compra", "Guía el checkout para reducir el abandono de carrito."),
         ("Soporte 24/7", "Disponible siempre, con escalado a un humano."),
         ("Mide impacto", "CSAT, intenciones y tasa de resolución en un dashboard.")]
for i, (t, d) in enumerate(feats):
    x = MX + (i % 2) * Inches(5.95); y = 3.35 + (i // 2) * 1.55
    card(s, x, Inches(y), Inches(5.6), Inches(1.3))
    rect(s, x + Inches(0.3), Inches(y + 0.32), Inches(0.18), Inches(0.18), ACCENT)
    textbox(s, x + Inches(0.65), Inches(y + 0.2), Inches(4.7), Inches(1.0),
            [{"text": t, "size": 15.5, "color": INK, "font": FSB, "bold": True, "sa": 4},
             {"text": d, "size": 12.5, "color": MUTED, "line": 1.15}])
footer(s)

# ===== 7 · SOLUCIÓN · DECISIONES: por qué no Azure =====
s = blank(); header(s, "Solución · Decisiones", "Decisión tecnológica: Azure y alternativas")
textbox(s, MX, Inches(2.0), Inches(11.4), Inches(0.85),
        [{"text": "El curso se centra en Azure AI; la opción ideal era Azure AI Language (CLU). Al agotarse el crédito de estudiante de Azure, se optó por alternativas de capa gratuita equivalentes, conservando una arquitectura desacoplada y migrable a Azure.", "size": 15, "color": MUTED, "line": 1.25}])
maps = [("Motor NLP", "Azure AI Language (CLU)", "LLM Groq (Llama 3.3)"),
        ("Backend", "Azure Container Instances", "Render"),
        ("Base de datos", "Azure DB for PostgreSQL", "PostgreSQL en Neon"),
        ("Frontend", "Azure Static Web Apps", "Netlify")]
textbox(s, MX + Inches(3.5), Inches(3.05), Inches(3.5), Inches(0.3), [{"text": "OPCIÓN IDEAL (AZURE)", "size": 10.5, "color": MUTED, "font": FSB, "bold": True}])
textbox(s, MX + Inches(7.7), Inches(3.05), Inches(3.5), Inches(0.3), [{"text": "IMPLEMENTADO (GRATIS)", "size": 10.5, "color": GREEN, "font": FSB, "bold": True}])
y = 3.4
for comp, az, alt in maps:
    card(s, MX, Inches(y), Inches(3.3), Inches(0.66), fill=LIGHT)
    textbox(s, MX + Inches(0.2), Inches(y), Inches(3.0), Inches(0.66), [{"text": comp, "size": 12.5, "color": ACCENT, "font": FSB, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, MX + Inches(3.5), Inches(y), Inches(3.6), Inches(0.66), [{"text": az, "size": 12.5, "color": MUTED}], anchor=MSO_ANCHOR.MIDDLE)
    arrow(s, MX + Inches(7.15), Inches(y + 0.33), MX + Inches(7.6), Inches(y + 0.33))
    textbox(s, MX + Inches(7.7), Inches(y), Inches(3.7), Inches(0.66), [{"text": alt, "size": 12.5, "color": INK, "font": FSB, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.76
textbox(s, MX, Inches(6.5), Inches(11.4), Inches(0.4),
        [{"text": "El motor de IA y la BD se cambian por variables de entorno: se puede volver a Azure sin reescribir el código.", "size": 12, "color": MUTED, "italic": True}])
footer(s)

# ===== 8 · SOLUCIÓN · TECNOLOGÍA: stack =====
s = blank(); header(s, "Solución · Tecnología", "Stack y justificación")
rows = [("Frontend", "HTML/CSS/JS vanilla", "Embebible con un <script>, sin dependencias; compatible con cualquier web."),
        ("Backend", "Python + FastAPI", "API REST rápida y asíncrona; documentación automática; ecosistema de IA."),
        ("Motor NLP", "LLM Groq (Llama 3.3)", "Entiende lenguaje libre con las FAQs como contexto; intercambiable con Azure CLU."),
        ("Base de datos", "PostgreSQL (Neon)", "Robusta y open source; ideal para logs y métricas; migrable a Azure DB."),
        ("Despliegue", "Render + Netlify", "Hosting gratuito; el backend resguarda las claves de API (no se exponen).")]
y = 2.1
for comp, tech, why in rows:
    card(s, MX, Inches(y), Inches(11.5), Inches(0.92))
    textbox(s, MX + Inches(0.3), Inches(y), Inches(2.3), Inches(0.92), [{"text": comp, "size": 13.5, "color": ACCENT, "font": FSB, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, MX + Inches(2.7), Inches(y), Inches(2.8), Inches(0.92), [{"text": tech, "size": 13.5, "color": INK, "font": FSB, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, MX + Inches(5.6), Inches(y), Inches(5.7), Inches(0.92), [{"text": why, "size": 12, "color": MUTED, "line": 1.15}], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.99
footer(s)

# ===== 9 · SOLUCIÓN · COSTOS =====
s = blank(); header(s, "Solución · Costos", "Costos de operación")
costs = [("Motor NLP", "Groq — Free", "Azure AI Language ≈ $25"),
         ("Backend / API", "Render — Free", "Azure App Service (B1) ≈ $13"),
         ("Base de datos", "Neon — Free", "Azure DB PostgreSQL ≈ $20"),
         ("Frontend / Hosting", "Netlify — Free", "Azure Static Web Apps ≈ $0–9"),
         ("Dominio + SSL", "Incluido — Free", "≈ $2")]
textbox(s, MX + Inches(4.0), Inches(2.0), Inches(3.5), Inches(0.3), [{"text": "MVP ACTUAL (CAPA GRATUITA)", "size": 10.5, "color": GREEN, "font": FSB, "bold": True}])
textbox(s, MX + Inches(7.8), Inches(2.0), Inches(3.6), Inches(0.3), [{"text": "PRODUCCIÓN EN AZURE (ESTIMADO/MES)", "size": 10.5, "color": MUTED, "font": FSB, "bold": True}])
y = 2.4
for comp, free, prod in costs:
    card(s, MX, Inches(y), Inches(11.5), Inches(0.62), fill=LIGHT)
    textbox(s, MX + Inches(0.25), Inches(y), Inches(3.6), Inches(0.62), [{"text": comp, "size": 12.5, "color": INK, "font": FSB, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, MX + Inches(4.0), Inches(y), Inches(3.5), Inches(0.62), [{"text": free, "size": 12.5, "color": GREEN}], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, MX + Inches(7.8), Inches(y), Inches(3.5), Inches(0.62), [{"text": prod, "size": 12.5, "color": MUTED}], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.7
card(s, MX, Inches(y), Inches(11.5), Inches(0.7), fill=ACCENT)
textbox(s, MX + Inches(0.25), Inches(y), Inches(3.6), Inches(0.7), [{"text": "TOTAL MENSUAL", "size": 13, "color": WHITE, "font": FSB, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, MX + Inches(4.0), Inches(y), Inches(3.5), Inches(0.7), [{"text": "US$ 0 / mes", "size": 15, "color": WHITE, "font": FSB, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, MX + Inches(7.8), Inches(y), Inches(3.5), Inches(0.7), [{"text": "≈ US$ 60 / mes (estimado)", "size": 13.5, "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, MX, Inches(6.55), Inches(11.4), Inches(0.35),
        [{"text": "Desarrollo: equipo de 3 personas, ~4 semanas (costo académico). Cifras de Azure referenciales según volumen.", "size": 11, "color": MUTED, "italic": True}])
footer(s)

# ===== 10 · SOLUCIÓN · DISEÑO: FAQs =====
s = blank(); header(s, "Solución · Diseño", "Base de conocimiento: FAQs por categoría")
cats = [("Productos", "Especificaciones, stock y garantía"), ("Precios", "Precios con IGV, descuentos y cupones"),
        ("Métodos de pago", "Yape, Plin, tarjetas y cuotas"), ("Proceso de compra", "Pasos, compra invitado, cancelación"),
        ("Envíos", "Tiempos, costos y rastreo de pedidos"), ("Soporte", "Garantías, devoluciones y post-venta")]
for i, (t, d) in enumerate(cats):
    x = MX + (i % 3) * Inches(3.85); y = 2.25 + (i // 3) * 1.75
    card(s, x, Inches(y), Inches(3.6), Inches(1.5))
    textbox(s, x + Inches(0.28), Inches(y + 0.22), Inches(3.1), Inches(1.1),
            [{"text": f"0{i+1}", "size": 13, "color": ACCENT, "font": FSB, "bold": True, "sa": 3},
             {"text": t, "size": 15.5, "color": INK, "font": FSB, "bold": True, "sa": 3},
             {"text": d, "size": 12, "color": MUTED, "line": 1.15}])
textbox(s, MX, Inches(6.0), Inches(11.4), Inches(0.5),
        [{"text": "Personalidad: «Nova», tono amigable y conciso, en español peruano. Umbral de confianza con derivación a humano (fallback).", "size": 12.5, "color": MUTED, "italic": True}])
footer(s)

# ===== 11 · SOLUCIÓN · DISEÑO: FLUJO =====
s = blank(); header(s, "Solución · Diseño", "Flujo conversacional")
steps = ["Usuario\nescribe", "Nova procesa\n(LLM + FAQs)", "Responde con\nla FAQ", "Registra en\nPostgreSQL", "Encuesta CSAT\ny cierre"]
bw, bh, gap = Inches(2.05), Inches(1.15), Inches(0.28)
x = MX; y = Inches(2.7)
for i, st in enumerate(steps):
    card(s, x, y, bw, bh, fill=LIGHT, border=LINE)
    textbox(s, x, y, bw, bh, [{"text": st, "size": 13, "color": INK, "font": FSB, "bold": True, "align": PP_ALIGN.CENTER, "line": 1.1}], anchor=MSO_ANCHOR.MIDDLE)
    if i < len(steps) - 1:
        ax = x + bw; arrow(s, ax, y + Inches(0.575), ax + gap, y + Inches(0.575))
    x = x + bw + gap
card(s, MX, Inches(4.45), Inches(11.5), Inches(0.95), fill=GREENB)
textbox(s, MX + Inches(0.3), Inches(4.45), Inches(11), Inches(0.95),
        [{"text": "Si la consulta no está cubierta por las FAQs → Nova lo reconoce (no inventa) y deriva a soporte / agente humano, registrando el caso como fallback para mejorar la base de conocimiento.", "size": 13.5, "color": INK, "line": 1.2}], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ===== 12 · SOLUCIÓN · DISEÑO: ARQUITECTURA =====
s = blank(); header(s, "Solución · Diseño", "Arquitectura técnica")
card(s, Inches(3.4), Inches(2.0), Inches(6.5), Inches(0.95), fill=LIGHT, border=LINE)
textbox(s, Inches(3.4), Inches(2.0), Inches(6.5), Inches(0.95), [{"text": "FRONTEND — Widget embebible + Tienda  (Netlify)", "size": 14, "color": INK, "font": FSB, "bold": True, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
arrow(s, Inches(6.65), Inches(2.95), Inches(6.65), Inches(3.35))
textbox(s, Inches(6.8), Inches(2.98), Inches(4), Inches(0.32), [{"text": "HTTPS · POST /api/chat", "size": 10.5, "color": MUTED, "italic": True}])
card(s, Inches(3.4), Inches(3.4), Inches(6.5), Inches(0.95), fill=ACCENT)
textbox(s, Inches(3.4), Inches(3.4), Inches(6.5), Inches(0.95), [{"text": "BACKEND — API REST FastAPI  (Render)", "size": 14, "color": WHITE, "font": FSB, "bold": True, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
arrow(s, Inches(5.2), Inches(4.35), Inches(4.4), Inches(4.85))
arrow(s, Inches(8.1), Inches(4.35), Inches(8.9), Inches(4.85))
card(s, Inches(1.7), Inches(4.9), Inches(4.6), Inches(1.15), fill=LIGHT, border=LINE)
textbox(s, Inches(1.7), Inches(4.9), Inches(4.6), Inches(1.15),
        [{"text": "MOTOR NLP", "size": 11, "color": ACCENT, "font": FSB, "bold": True, "align": PP_ALIGN.CENTER, "sa": 2},
         {"text": "LLM Groq (Llama 3.3) + FAQs\nAlternativa: Azure AI Language (CLU)", "size": 11.5, "color": INK, "align": PP_ALIGN.CENTER, "line": 1.1}], anchor=MSO_ANCHOR.MIDDLE)
card(s, Inches(7.0), Inches(4.9), Inches(4.6), Inches(1.15), fill=LIGHT, border=LINE)
textbox(s, Inches(7.0), Inches(4.9), Inches(4.6), Inches(1.15),
        [{"text": "BASE DE DATOS", "size": 11, "color": ACCENT, "font": FSB, "bold": True, "align": PP_ALIGN.CENTER, "sa": 2},
         {"text": "PostgreSQL (Neon)\nfaqs · chat_logs · sessions · feedback", "size": 11.5, "color": INK, "align": PP_ALIGN.CENTER, "line": 1.1}], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, MX, Inches(6.35), Inches(11.4), Inches(0.4),
        [{"text": "Motor NLP y base de datos son piezas intercambiables (se cambian por variables de entorno).", "size": 12, "color": MUTED, "italic": True, "align": PP_ALIGN.CENTER}])
footer(s)

# ===== 13 · IMPLEMENTACIÓN: PLAN =====
s = blank(); header(s, "Implementación", "Pasos para integrar el chatbot al sitio")
phases = [("1", "Diseño y FAQs", "Levantar FAQs con ventas; definir flujo e intenciones."),
          ("2", "Backend + BD", "FastAPI + PostgreSQL; endpoints y registro de logs."),
          ("3", "Bot + Frontend", "Conectar el LLM y construir el widget + CSAT."),
          ("4", "Integración", "Insertar el <script> en el sitio; configurar CORS."),
          ("5", "Despliegue", "Publicar en Render + Netlify; pruebas y monitoreo.")]
bw = Inches(2.12); gap = Inches(0.18); y = Inches(3.0); x = MX
rect(s, MX + Inches(0.1), Inches(2.92), Inches(11.2), Pt(2), LINE)
for num, t, d in phases:
    card(s, x, y, bw, Inches(1.7))
    textbox(s, x, y + Inches(0.18), bw, Inches(0.5), [{"text": num, "size": 26, "color": ACCENT, "font": FL, "align": PP_ALIGN.CENTER}])
    textbox(s, x + Inches(0.12), y + Inches(0.72), bw - Inches(0.24), Inches(0.9),
            [{"text": t, "size": 13, "color": INK, "font": FSB, "bold": True, "align": PP_ALIGN.CENTER, "sa": 3},
             {"text": d, "size": 10.5, "color": MUTED, "align": PP_ALIGN.CENTER, "line": 1.1}])
    x += bw + gap
footer(s)

# ===== 14 · IMPLEMENTACIÓN: MÉTRICAS =====
s = blank(); header(s, "Implementación", "Métricas de éxito (objetivos)")
kpis = [("≥ 85%", "Resolución sin agente"), ("≥ 4.2/5", "CSAT (satisfacción)"),
        ("−25%", "Carritos abandonados"), ("< 2 s", "Tiempo de respuesta")]
for i, (v, l) in enumerate(kpis):
    x = MX + i * Inches(2.92)
    card(s, x, Inches(2.2), Inches(2.7), Inches(1.7))
    textbox(s, x, Inches(2.45), Inches(2.7), Inches(0.8), [{"text": v, "size": 34, "color": ACCENT, "font": FL, "align": PP_ALIGN.CENTER}])
    textbox(s, x, Inches(3.25), Inches(2.7), Inches(0.5), [{"text": l, "size": 12.5, "color": MUTED, "align": PP_ALIGN.CENTER}])
textbox(s, MX, Inches(4.35), Inches(11.4), Inches(0.4), [{"text": "QUÉ SE MONITOREA (en vivo, desde PostgreSQL)", "size": 12, "color": ACCENT, "font": FSB, "bold": True}])
mons = ["Total de conversaciones", "Top intenciones consultadas", "Tasa de fallback (no entendidas)",
        "CSAT promedio (1–5)", "Tasa de resolución / escalado", "Conversiones atribuidas al chat"]
for i, m in enumerate(mons):
    x = MX + (i % 3) * Inches(3.85); y = 4.85 + (i // 3) * 0.7
    rect(s, x, Inches(y + 0.08), Inches(0.16), Inches(0.16), GREEN)
    textbox(s, x + Inches(0.32), Inches(y), Inches(3.4), Inches(0.5), [{"text": m, "size": 12.5, "color": INK}], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ===== 15-16 · PREGUNTAS GUÍA =====
def pregunta(slide, y, q, a):
    textbox(slide, MX, Inches(y), Inches(11.4), Inches(0.5), [{"text": q, "size": 15, "color": INK, "font": FSB, "bold": True}])
    textbox(slide, MX, Inches(y + 0.42), Inches(11.4), Inches(0.9), [{"text": a, "size": 13, "color": MUTED, "line": 1.2}])

s = blank(); header(s, "Preguntas guía", "Respuestas (1 – 3)")
pregunta(s, 2.05, "1) ¿Cuáles son las FAQs más comunes de los clientes?",
         "Especificaciones y disponibilidad de productos, garantía; precios, IGV y descuentos; métodos de pago (Yape, Plin, tarjetas, cuotas); proceso de compra y cancelación; tiempos/costos y rastreo de envíos; devoluciones y soporte técnico.")
pregunta(s, 3.6, "2) ¿Qué herramientas/plataformas son más adecuadas?",
         "Widget en JS embebible (Netlify) + API FastAPI (Render) + motor NLP con LLM Groq usando las FAQs como contexto (alineado al curso vía Azure AI Language) + PostgreSQL (Neon) para métricas. Gratis, escalable e intercambiable.")
pregunta(s, 5.15, "3) ¿Cómo evaluar la efectividad (satisfacción y ventas)?",
         "Satisfacción: CSAT (1–5), tasa de resolución sin escalado y tasa de comprensión. Ventas: reducción de carritos abandonados y conversión post-chat. Todo registrado en PostgreSQL y visible en el dashboard.")
footer(s)

s = blank(); header(s, "Preguntas guía", "Respuestas (4 – 5)")
pregunta(s, 2.15, "4) ¿Qué desafíos técnicos podrían surgir y cómo resolverlos?",
         "Comprensión del lenguaje → LLM con grounding en FAQs (no inventa). Integración → widget <script> + CORS. Seguridad → las claves viven en el backend, nunca en el navegador. Disponibilidad → contenedores con healthcheck/restart. FAQs actualizadas → edición en base de datos.")
pregunta(s, 4.2, "5) ¿Qué métricas monitorear y cómo optimizar las funcionalidades?",
         "Monitorear: nº de conversaciones, top intenciones, % de fallback, CSAT y tasa de resolución. Optimizar: revisar los mensajes en fallback para convertirlos en nuevas FAQs, ajustar el prompt/base de conocimiento y re-entrenar periódicamente con casos reales.")
footer(s)

# ===== 17 · DEMO EN VIVO =====
s = blank(); header(s, "Resultado", "Demo en vivo")
demo = [("Tienda + Widget", "Sitio de InnovVentas con Nova flotante (Netlify)"),
        ("Chatbot Nova", "Conversación real con persistencia del historial"),
        ("Dashboard", "Métricas en vivo desde la API")]
for i, (t, d) in enumerate(demo):
    x = MX + i * Inches(3.9)
    card(s, x, Inches(2.2), Inches(3.65), Inches(1.7))
    textbox(s, x + Inches(0.25), Inches(2.45), Inches(3.2), Inches(1.3),
            [{"text": t, "size": 15.5, "color": INK, "font": FSB, "bold": True, "sa": 5},
             {"text": d, "size": 12, "color": MUTED, "line": 1.2}])
card(s, MX, Inches(4.35), Inches(11.5), Inches(1.5), fill=LIGHT)
textbox(s, MX + Inches(0.35), Inches(4.55), Inches(11), Inches(1.1),
        [{"text": "URLs del demo", "size": 12, "color": ACCENT, "font": FSB, "bold": True, "sa": 6},
         {"runs": [{"text": "API (backend):  ", "size": 13.5, "color": INK, "font": FSB, "bold": True},
                   {"text": "innovventas-chatbot-rz1k.onrender.com", "size": 13.5, "color": ACCENT}], "sa": 4},
         {"runs": [{"text": "Tienda + Dashboard:  ", "size": 13.5, "color": INK, "font": FSB, "bold": True},
                   {"text": "publicado en Netlify (frontend estático)", "size": 13.5, "color": MUTED}]}])
footer(s)

# ===== 18 · CONCLUSIONES =====
s = blank(); header(s, "Cierre", "Conclusiones")
concl = ["Nova ataca la causa raíz del abandono de carrito: atención inmediata, 24/7 e integrada al sitio.",
         "Solución funcional y desplegada en la nube con costo cero (capa gratuita).",
         "Arquitectura desacoplada: el motor de IA es intercambiable (Groq hoy, Azure AI Language en producción).",
         "Cada conversación se registra y se mide, habilitando optimización continua basada en datos."]
y = 2.2
for c in concl:
    rect(s, MX, Inches(y + 0.06), Inches(0.16), Inches(0.16), ACCENT)
    textbox(s, MX + Inches(0.4), Inches(y - 0.05), Inches(10.9), Inches(0.9), [{"text": c, "size": 15, "color": INK, "line": 1.2}])
    y += 1.0
footer(s)

# ===== 19 · GRACIAS =====
s = blank()
rect(s, 0, 0, Inches(0.22), Inches(7.5), ACCENT)
textbox(s, MX, Inches(2.6), Inches(11.5), Inches(1.5), [{"text": "¡Gracias!", "size": 60, "color": INK, "font": FL}])
textbox(s, MX, Inches(4.0), Inches(11.5), Inches(0.6), [{"text": "¿Preguntas? Prueba a Nova en vivo en el sitio de InnovVentas.", "size": 16, "color": MUTED}])
rect(s, MX, Inches(4.75), Inches(2.2), Pt(2.5), LINE)
textbox(s, MX, Inches(5.0), Inches(11.5), Inches(0.9),
        [{"runs": [{"text": "soporte@innovventas.pe", "size": 13.5, "color": INK, "font": FSB, "bold": True},
                   {"text": "    ·    01-234-5678    ·    ", "size": 13.5, "color": MUTED},
                   {"text": "innovventas-chatbot-rz1k.onrender.com", "size": 13.5, "color": ACCENT}]}])

prs.core_properties.title = "InnovVentas — Chatbot Nova | Trabajo Final AI-900T00"
prs.core_properties.author = "Cahua · Cruces · Olivitos — SENATI"
prs.save(OUTPUT)
chk = Presentation(OUTPUT)
print(f"OK -> {OUTPUT}  ({len(chk.slides)} slides)")
